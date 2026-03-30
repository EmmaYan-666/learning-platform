#!/usr/bin/env python3
"""
系统数据同步问题分析PPT生成器 - 优化版
改进布局、对齐和视觉层次
"""

import subprocess
import sys
import os

def check_and_install_dependencies():
    """检查并安装所需的Python库"""
    print("检查依赖库...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        print("✓ python-pptx已安装")
        return True
    except ImportError:
        print("× python-pptx未安装,正在安装...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            print("✓ python-pptx安装成功")
            return True
        except subprocess.CalledProcessError as e:
            print(f"✗ 安装失败: {e}")
            print("\n请手动安装: pip3 install python-pptx")
            return False

def create_presentation():
    """创建系统问题分析PPT - 优化版"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 添加空白幻灯片
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # 颜色定义 - 专业配色方案
    colors = {
        'primary': RGBColor(30, 39, 97),        # 深蓝色
        'secondary': RGBColor(202, 220, 252),   # 浅蓝色
        'accent': RGBColor(249, 97, 103),       # 珊瑚红
        'text': RGBColor(44, 62, 80),           # 深灰色
        'light_gray': RGBColor(248, 249, 250),  # 浅灰色
        'success': RGBColor(44, 95, 45),        # 森林绿
        'warning': RGBColor(243, 156, 18),      # 橙色
        'white': RGBColor(255, 255, 255),
    }
    
    # ========== 标题 ==========
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "系统数据同步问题分析与解决方案"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors['primary']
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.margin_left = 0
    title_frame.margin_right = 0
    title_frame.margin_top = 0
    title_frame.margin_bottom = 0
    
    # ========== 左列: 系统架构 ==========
    # 标题背景条
    section1_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(0.06), Inches(0.28)
    )
    section1_bar.fill.solid()
    section1_bar.fill.fore_color.rgb = colors['primary']
    section1_bar.line.fill.background()
    
    # 标题
    section1_title = slide.shapes.add_textbox(Inches(0.65), Inches(0.85), Inches(2), Inches(0.28))
    section1_frame = section1_title.text_frame
    section1_frame.text = "系统架构"
    section1_frame.paragraphs[0].font.size = Pt(14)
    section1_frame.paragraphs[0].font.bold = True
    section1_frame.paragraphs[0].font.color.rgb = colors['primary']
    section1_frame.paragraphs[0].font.name = "Arial"
    section1_frame.margin_left = 0
    section1_frame.margin_top = 0
    
    # 系统架构图 - 水平排列三个框
    box_width = Inches(1.4)
    box_height = Inches(0.65)
    box_y = Inches(1.25)
    
    # B360系统框
    b360_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), box_y, box_width, box_height
    )
    b360_shape.fill.solid()
    b360_shape.fill.fore_color.rgb = colors['primary']
    b360_shape.line.fill.background()
    
    b360_tf = b360_shape.text_frame
    b360_tf.text = "B360\n中央预定系统"
    b360_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    b360_tf.paragraphs[0].font.size = Pt(10)
    b360_tf.paragraphs[0].font.bold = True
    b360_tf.paragraphs[0].font.color.rgb = colors['white']
    b360_tf.paragraphs[0].font.name = "Arial"
    b360_tf.word_wrap = True
    
    # 箭头1
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(1.95), Inches(1.42), Inches(0.35), Inches(0.25)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = colors['text']
    arrow1.line.fill.background()
    
    # OXI接口框
    oxi_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), box_y, Inches(1.0), box_height
    )
    oxi_shape.fill.solid()
    oxi_shape.fill.fore_color.rgb = colors['secondary']
    oxi_shape.line.fill.background()
    
    oxi_tf = oxi_shape.text_frame
    oxi_tf.text = "OXI接口"
    oxi_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    oxi_tf.paragraphs[0].font.size = Pt(10)
    oxi_tf.paragraphs[0].font.bold = True
    oxi_tf.paragraphs[0].font.color.rgb = colors['primary']
    oxi_tf.paragraphs[0].font.name = "Arial"
    
    # 箭头2
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(1.42), Inches(0.35), Inches(0.25)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = colors['text']
    arrow2.line.fill.background()
    
    # PMS系统框
    pms_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), box_y, box_width, box_height
    )
    pms_shape.fill.solid()
    pms_shape.fill.fore_color.rgb = colors['primary']
    pms_shape.line.fill.background()
    
    pms_tf = pms_shape.text_frame
    pms_tf.text = "PMS\n酒店管理系统"
    pms_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    pms_tf.paragraphs[0].font.size = Pt(10)
    pms_tf.paragraphs[0].font.bold = True
    pms_tf.paragraphs[0].font.color.rgb = colors['white']
    pms_tf.paragraphs[0].font.name = "Arial"
    pms_tf.word_wrap = True
    
    # 同步机制说明框
    sync_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.05), Inches(4.7), Inches(0.35)
    )
    sync_box.fill.solid()
    sync_box.fill.fore_color.rgb = RGBColor(240, 243, 247)
    sync_box.line.color.rgb = RGBColor(220, 225, 230)
    sync_box.line.width = Pt(1)
    
    sync_tf = sync_box.text_frame
    sync_tf.text = "同步机制: 异步消息队列 + 全量覆盖"
    sync_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sync_tf.paragraphs[0].font.size = Pt(9)
    sync_tf.paragraphs[0].font.color.rgb = colors['text']
    sync_tf.paragraphs[0].font.name = "Arial"
    
    # ========== 左列: 业务需求 ==========
    # 标题背景条
    section2_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.55), Inches(0.06), Inches(0.28)
    )
    section2_bar.fill.solid()
    section2_bar.fill.fore_color.rgb = colors['warning']
    section2_bar.line.fill.background()
    
    # 标题
    section2_title = slide.shapes.add_textbox(Inches(0.65), Inches(2.55), Inches(2.5), Inches(0.28))
    section2_frame = section2_title.text_frame
    section2_frame.text = "业务需求"
    section2_frame.paragraphs[0].font.size = Pt(14)
    section2_frame.paragraphs[0].font.bold = True
    section2_frame.paragraphs[0].font.color.rgb = colors['warning']
    section2_frame.paragraphs[0].font.name = "Arial"
    section2_frame.margin_left = 0
    section2_frame.margin_top = 0
    
    # 业务需求内容框
    req_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.95), Inches(4.7), Inches(0.75)
    )
    req_box.fill.solid()
    req_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    req_box.line.color.rgb = colors['warning']
    req_box.line.width = Pt(1)
    
    req_tf = req_box.text_frame
    req_tf.word_wrap = True
    req_tf.text = "当客人的会员SLC (Welcome Amenity) 满足条件时,\n自动往订单添加special code代表此次入住给客人的权益"
    req_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    req_tf.paragraphs[0].font.size = Pt(9)
    req_tf.paragraphs[0].font.color.rgb = colors['text']
    req_tf.paragraphs[0].font.name = "Arial"
    req_tf.margin_left = Pt(8)
    req_tf.margin_right = Pt(8)
    req_tf.margin_top = Pt(6)
    
    # ========== 右列: 核心问题 ==========
    # 标题背景条
    section3_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(0.85), Inches(0.06), Inches(0.28)
    )
    section3_bar.fill.solid()
    section3_bar.fill.fore_color.rgb = colors['accent']
    section3_bar.line.fill.background()
    
    # 标题
    section3_title = slide.shapes.add_textbox(Inches(5.45), Inches(0.85), Inches(2.5), Inches(0.28))
    section3_frame = section3_title.text_frame
    section3_frame.text = "核心问题"
    section3_frame.paragraphs[0].font.size = Pt(14)
    section3_frame.paragraphs[0].font.bold = True
    section3_frame.paragraphs[0].font.color.rgb = colors['accent']
    section3_frame.paragraphs[0].font.name = "Arial"
    section3_frame.margin_left = 0
    section3_frame.margin_top = 0
    
    # 核心问题内容框
    prob_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.9)
    )
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
    prob_box.line.color.rgb = colors['accent']
    prob_box.line.width = Pt(1)
    
    # 问题列表
    problems = [
        "❌ B360处理SLC逻辑,订单双向同步",
        "❌ PMS FO修改订单时,B360推送更新可能失败",
        "❌ 全量覆盖导致PMS新修改内容丢失",
        "    例如:重要的Comments记录被覆盖"
    ]
    
    prob_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.35), Inches(4.0), Inches(1.7))
    prob_tf = prob_text.text_frame
    prob_tf.word_wrap = True
    
    for i, prob in enumerate(problems):
        if i == 0:
            p = prob_tf.paragraphs[0]
        else:
            p = prob_tf.add_paragraph()
        p.text = prob
        p.font.size = Pt(9)
        p.font.color.rgb = colors['text']
        p.font.name = "Arial"
        p.space_before = Pt(6)
        p.space_after = Pt(0)
    
    # ========== 中间: 问题流程示意 ==========
    # 标题背景条
    section4_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.85), Inches(0.06), Inches(0.28)
    )
    section4_bar.fill.solid()
    section4_bar.fill.fore_color.rgb = colors['accent']
    section4_bar.line.fill.background()
    
    # 标题
    section4_title = slide.shapes.add_textbox(Inches(0.65), Inches(3.85), Inches(2.5), Inches(0.28))
    section4_frame = section4_title.text_frame
    section4_frame.text = "问题流程示意"
    section4_frame.paragraphs[0].font.size = Pt(14)
    section4_frame.paragraphs[0].font.bold = True
    section4_frame.paragraphs[0].font.color.rgb = colors['accent']
    section4_frame.paragraphs[0].font.name = "Arial"
    section4_frame.margin_left = 0
    section4_frame.margin_top = 0
    
    # 流程步骤 - 4个步骤水平排列
    step_width = Inches(2.0)
    step_height = Inches(0.55)
    step_y = Inches(4.25)
    arrow_width = Inches(0.25)
    
    # 步骤1
    step1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), step_y, step_width, step_height
    )
    step1.fill.solid()
    step1.fill.fore_color.rgb = colors['light_gray']
    step1.line.color.rgb = colors['text']
    step1.line.width = Pt(1)
    
    step1_tf = step1.text_frame
    step1_tf.text = "① PMS FO修改订单\n    (添加Comments)"
    step1_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step1_tf.paragraphs[0].font.size = Pt(9)
    step1_tf.paragraphs[0].font.color.rgb = colors['text']
    step1_tf.paragraphs[0].font.name = "Arial"
    step1_tf.margin_left = Pt(8)
    step1_tf.margin_top = Pt(4)
    
    # 箭头
    flow_arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(2.55), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow1.fill.solid()
    flow_arrow1.fill.fore_color.rgb = colors['warning']
    flow_arrow1.line.fill.background()
    
    # 步骤2
    step2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.85), step_y, step_width, step_height
    )
    step2.fill.solid()
    step2.fill.fore_color.rgb = colors['light_gray']
    step2.line.color.rgb = colors['text']
    step2.line.width = Pt(1)
    
    step2_tf = step2.text_frame
    step2_tf.text = "② 同步至B360\n    触发SLC逻辑"
    step2_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step2_tf.paragraphs[0].font.size = Pt(9)
    step2_tf.paragraphs[0].font.color.rgb = colors['text']
    step2_tf.paragraphs[0].font.name = "Arial"
    step2_tf.margin_left = Pt(8)
    step2_tf.margin_top = Pt(4)
    
    # 箭头
    flow_arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.9), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow2.fill.solid()
    flow_arrow2.fill.fore_color.rgb = colors['accent']
    flow_arrow2.line.fill.background()
    
    # 步骤3 - 高亮
    step3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), step_y, step_width, step_height
    )
    step3.fill.solid()
    step3.fill.fore_color.rgb = RGBColor(255, 235, 235)
    step3.line.color.rgb = colors['accent']
    step3.line.width = Pt(2)
    
    step3_tf = step3.text_frame
    step3_tf.text = "③ B360推回更新\n    (删除special code)"
    step3_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step3_tf.paragraphs[0].font.size = Pt(9)
    step3_tf.paragraphs[0].font.bold = True
    step3_tf.paragraphs[0].font.color.rgb = colors['accent']
    step3_tf.paragraphs[0].font.name = "Arial"
    step3_tf.margin_left = Pt(8)
    step3_tf.margin_top = Pt(4)
    
    # 箭头
    flow_arrow3 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(7.25), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow3.fill.solid()
    flow_arrow3.fill.fore_color.rgb = colors['accent']
    flow_arrow3.line.fill.background()
    
    # 步骤4 - 问题结果
    step4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.55), step_y, Inches(1.95), step_height
    )
    step4.fill.solid()
    step4.fill.fore_color.rgb = colors['accent']
    step4.line.fill.background()
    
    step4_tf = step4.text_frame
    step4_tf.text = "❌ 更新失败\n或覆盖新数据"
    step4_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    step4_tf.paragraphs[0].font.size = Pt(9)
    step4_tf.paragraphs[0].font.bold = True
    step4_tf.paragraphs[0].font.color.rgb = colors['white']
    step4_tf.paragraphs[0].font.name = "Arial"
    
    # ========== 底部: 解决方案 ==========
    # 标题背景条
    section5_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.95), Inches(0.06), Inches(0.28)
    )
    section5_bar.fill.solid()
    section5_bar.fill.fore_color.rgb = colors['success']
    section5_bar.line.fill.background()
    
    # 标题
    section5_title = slide.shapes.add_textbox(Inches(0.65), Inches(4.95), Inches(2.5), Inches(0.28))
    section5_frame = section5_title.text_frame
    section5_frame.text = "建议解决方案"
    section5_frame.paragraphs[0].font.size = Pt(14)
    section5_frame.paragraphs[0].font.bold = True
    section5_frame.paragraphs[0].font.color.rgb = colors['success']
    section5_frame.paragraphs[0].font.name = "Arial"
    section5_frame.margin_left = 0
    section5_frame.margin_top = 0
    
    # 解决方案框
    solution_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.35), Inches(9.0), Inches(0.55)
    )
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    solution_box.line.color.rgb = colors['success']
    solution_box.line.width = Pt(2)
    
    # 解决方案内容
    solution_items = [
        ("✓", "将SLC权益逻辑迁移至PMS端处理"),
        ("→", "PMS本地监控订单变更,自动决定special code更新"),
        ("→", "所有变更由PMS统一发起"),
        ("→", "B360仅作为数据接收方"),
        ("→", "避免双向覆盖冲突")
    ]
    
    solution_text = slide.shapes.add_textbox(Inches(0.6), Inches(5.35), Inches(8.8), Inches(0.55))
    solution_tf = solution_text.text_frame
    solution_tf.word_wrap = True
    
    for i, (symbol, text) in enumerate(solution_items):
        if i == 0:
            p = solution_tf.paragraphs[0]
        else:
            p = solution_tf.add_paragraph()
        p.text = f"{symbol} {text}"
        p.font.size = Pt(10)
        p.font.color.rgb = colors['text']
        p.font.name = "Arial"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
        if i == 0:
            p.font.bold = True
    
    # 保存文件
    output_path = "/Users/emma/WorkBuddy/20260312092945/系统数据同步问题.pptx"
    prs.save(output_path)
    print(f"\n✓ PPT已成功创建: {output_path}")
    return output_path

def main():
    """主函数"""
    print("=" * 60)
    print("系统数据同步问题分析PPT生成器 - 优化版")
    print("=" * 60)
    print()
    
    # 检查并安装依赖
    if not check_and_install_dependencies():
        sys.exit(1)
    
    # 创建演示文稿
    output_file = create_presentation()
    
    print("\n" + "=" * 60)
    print("完成! PPT已优化布局,更加清晰易读")
    print("=" * 60)

if __name__ == "__main__":
    main()
