#!/usr/bin/env python3
"""
System Data Synchronization Issue Analysis PPT Generator - English Version
"""

import subprocess
import sys
import os

def check_and_install_dependencies():
    """Check and install required Python libraries"""
    print("Checking dependencies...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        print("✓ python-pptx installed")
        return True
    except ImportError:
        print("× python-pptx not installed, installing...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            print("✓ python-pptx installed successfully")
            return True
        except subprocess.CalledProcessError as e:
            print(f"✗ Installation failed: {e}")
            print("\nPlease install manually: pip3 install python-pptx")
            return False

def create_presentation():
    """Create System Issue Analysis PPT - English Version"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Color scheme
    colors = {
        'primary': RGBColor(30, 39, 97),        # Navy blue
        'secondary': RGBColor(202, 220, 252),   # Light blue
        'accent': RGBColor(249, 97, 103),       # Coral red
        'text': RGBColor(44, 62, 80),           # Dark gray
        'light_gray': RGBColor(248, 249, 250),  # Light gray
        'success': RGBColor(44, 95, 45),        # Forest green
        'warning': RGBColor(243, 156, 18),      # Orange
        'white': RGBColor(255, 255, 255),
    }
    
    # ========== Title ==========
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "System Data Synchronization Issue Analysis & Solution"
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = colors['primary']
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.margin_left = 0
    title_frame.margin_right = 0
    title_frame.margin_top = 0
    title_frame.margin_bottom = 0
    
    # ========== Left Column: System Architecture ==========
    # Section title bar
    section1_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.85), Inches(0.06), Inches(0.28)
    )
    section1_bar.fill.solid()
    section1_bar.fill.fore_color.rgb = colors['primary']
    section1_bar.line.fill.background()
    
    # Section title
    section1_title = slide.shapes.add_textbox(Inches(0.65), Inches(0.85), Inches(2), Inches(0.28))
    section1_frame = section1_title.text_frame
    section1_frame.text = "System Architecture"
    section1_frame.paragraphs[0].font.size = Pt(14)
    section1_frame.paragraphs[0].font.bold = True
    section1_frame.paragraphs[0].font.color.rgb = colors['primary']
    section1_frame.paragraphs[0].font.name = "Arial"
    section1_frame.margin_left = 0
    section1_frame.margin_top = 0
    
    # System architecture diagram - 3 boxes horizontally
    box_width = Inches(1.4)
    box_height = Inches(0.65)
    box_y = Inches(1.25)
    
    # B360 box
    b360_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), box_y, box_width, box_height
    )
    b360_shape.fill.solid()
    b360_shape.fill.fore_color.rgb = colors['primary']
    b360_shape.line.fill.background()
    
    b360_tf = b360_shape.text_frame
    b360_tf.text = "B360\nCentral Reservation"
    b360_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    b360_tf.paragraphs[0].font.size = Pt(10)
    b360_tf.paragraphs[0].font.bold = True
    b360_tf.paragraphs[0].font.color.rgb = colors['white']
    b360_tf.paragraphs[0].font.name = "Arial"
    b360_tf.word_wrap = True
    
    # Arrow 1
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(1.95), Inches(1.42), Inches(0.35), Inches(0.25)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = colors['text']
    arrow1.line.fill.background()
    
    # OXI interface box
    oxi_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.35), box_y, Inches(1.0), box_height
    )
    oxi_shape.fill.solid()
    oxi_shape.fill.fore_color.rgb = colors['secondary']
    oxi_shape.line.fill.background()
    
    oxi_tf = oxi_shape.text_frame
    oxi_tf.text = "OXI\nInterface"
    oxi_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    oxi_tf.paragraphs[0].font.size = Pt(10)
    oxi_tf.paragraphs[0].font.bold = True
    oxi_tf.paragraphs[0].font.color.rgb = colors['primary']
    oxi_tf.paragraphs[0].font.name = "Arial"
    
    # Arrow 2
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(3.4), Inches(1.42), Inches(0.35), Inches(0.25)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = colors['text']
    arrow2.line.fill.background()
    
    # PMS box
    pms_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), box_y, box_width, box_height
    )
    pms_shape.fill.solid()
    pms_shape.fill.fore_color.rgb = colors['primary']
    pms_shape.line.fill.background()
    
    pms_tf = pms_shape.text_frame
    pms_tf.text = "PMS\nProperty Mgmt"
    pms_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    pms_tf.paragraphs[0].font.size = Pt(10)
    pms_tf.paragraphs[0].font.bold = True
    pms_tf.paragraphs[0].font.color.rgb = colors['white']
    pms_tf.paragraphs[0].font.name = "Arial"
    pms_tf.word_wrap = True
    
    # Sync mechanism box
    sync_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.05), Inches(4.7), Inches(0.35)
    )
    sync_box.fill.solid()
    sync_box.fill.fore_color.rgb = RGBColor(240, 243, 247)
    sync_box.line.color.rgb = RGBColor(220, 225, 230)
    sync_box.line.width = Pt(1)
    
    sync_tf = sync_box.text_frame
    sync_tf.text = "Sync Mechanism: Async Message Queue + Full Override"
    sync_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sync_tf.paragraphs[0].font.size = Pt(9)
    sync_tf.paragraphs[0].font.color.rgb = colors['text']
    sync_tf.paragraphs[0].font.name = "Arial"
    
    # ========== Left Column: Business Requirement ==========
    # Section title bar
    section2_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.55), Inches(0.06), Inches(0.28)
    )
    section2_bar.fill.solid()
    section2_bar.fill.fore_color.rgb = colors['warning']
    section2_bar.line.fill.background()
    
    # Section title
    section2_title = slide.shapes.add_textbox(Inches(0.65), Inches(2.55), Inches(2.5), Inches(0.28))
    section2_frame = section2_title.text_frame
    section2_frame.text = "Business Requirement"
    section2_frame.paragraphs[0].font.size = Pt(14)
    section2_frame.paragraphs[0].font.bold = True
    section2_frame.paragraphs[0].font.color.rgb = colors['warning']
    section2_frame.paragraphs[0].font.name = "Arial"
    section2_frame.margin_left = 0
    section2_frame.margin_top = 0
    
    # Business requirement content box
    req_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.95), Inches(4.7), Inches(0.75)
    )
    req_box.fill.solid()
    req_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    req_box.line.color.rgb = colors['warning']
    req_box.line.width = Pt(1)
    
    req_tf = req_box.text_frame
    req_tf.word_wrap = True
    req_tf.text = "When member's SLC (Welcome Amenity) meets criteria,\nautomatically add special code to reservation for guest benefits"
    req_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    req_tf.paragraphs[0].font.size = Pt(9)
    req_tf.paragraphs[0].font.color.rgb = colors['text']
    req_tf.paragraphs[0].font.name = "Arial"
    req_tf.margin_left = Pt(8)
    req_tf.margin_right = Pt(8)
    req_tf.margin_top = Pt(6)
    
    # ========== Right Column: Core Issues ==========
    # Section title bar
    section3_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(0.85), Inches(0.06), Inches(0.28)
    )
    section3_bar.fill.solid()
    section3_bar.fill.fore_color.rgb = colors['accent']
    section3_bar.line.fill.background()
    
    # Section title
    section3_title = slide.shapes.add_textbox(Inches(5.45), Inches(0.85), Inches(2.5), Inches(0.28))
    section3_frame = section3_title.text_frame
    section3_frame.text = "Core Issues"
    section3_frame.paragraphs[0].font.size = Pt(14)
    section3_frame.paragraphs[0].font.bold = True
    section3_frame.paragraphs[0].font.color.rgb = colors['accent']
    section3_frame.paragraphs[0].font.name = "Arial"
    section3_frame.margin_left = 0
    section3_frame.margin_top = 0
    
    # Core issues content box
    prob_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.25), Inches(4.2), Inches(1.9)
    )
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = RGBColor(255, 245, 245)
    prob_box.line.color.rgb = colors['accent']
    prob_box.line.width = Pt(1)
    
    # Issue list
    problems = [
        "❌ B360 handles SLC logic, causing bi-directional sync",
        "❌ When PMS FO modifies reservation, B360 push may fail",
        "❌ Full override causes PMS new changes to be lost",
        "    Example: Important Comments records overwritten"
    ]
    
    prob_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.35), Inches(4.0), Inches(1.7))
    prob_tf = prob_text.text_frame
    prob_tf.word_wrap = True
    
    for i, prob in enumerate(problems):
        if i == 0:
            p = prob_tf.paragraphs[0]
        else:
            p = prob_tf.add_paragraph()
        p.text = prob
        p.font.size = Pt(9)
        p.font.color.rgb = colors['text']
        p.font.name = "Arial"
        p.space_before = Pt(6)
        p.space_after = Pt(0)
    
    # ========== Middle: Problem Flow ==========
    # Section title bar
    section4_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.85), Inches(0.06), Inches(0.28)
    )
    section4_bar.fill.solid()
    section4_bar.fill.fore_color.rgb = colors['accent']
    section4_bar.line.fill.background()
    
    # Section title
    section4_title = slide.shapes.add_textbox(Inches(0.65), Inches(3.85), Inches(2.5), Inches(0.28))
    section4_frame = section4_title.text_frame
    section4_frame.text = "Issue Flow Diagram"
    section4_frame.paragraphs[0].font.size = Pt(14)
    section4_frame.paragraphs[0].font.bold = True
    section4_frame.paragraphs[0].font.color.rgb = colors['accent']
    section4_frame.paragraphs[0].font.name = "Arial"
    section4_frame.margin_left = 0
    section4_frame.margin_top = 0
    
    # Flow steps - 4 steps horizontally
    step_width = Inches(2.0)
    step_height = Inches(0.55)
    step_y = Inches(4.25)
    arrow_width = Inches(0.25)
    
    # Step 1
    step1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), step_y, step_width, step_height
    )
    step1.fill.solid()
    step1.fill.fore_color.rgb = colors['light_gray']
    step1.line.color.rgb = colors['text']
    step1.line.width = Pt(1)
    
    step1_tf = step1.text_frame
    step1_tf.text = "① PMS FO modifies\n    reservation (Comments)"
    step1_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step1_tf.paragraphs[0].font.size = Pt(9)
    step1_tf.paragraphs[0].font.color.rgb = colors['text']
    step1_tf.paragraphs[0].font.name = "Arial"
    step1_tf.margin_left = Pt(8)
    step1_tf.margin_top = Pt(4)
    
    # Arrow
    flow_arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(2.55), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow1.fill.solid()
    flow_arrow1.fill.fore_color.rgb = colors['warning']
    flow_arrow1.line.fill.background()
    
    # Step 2
    step2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.85), step_y, step_width, step_height
    )
    step2.fill.solid()
    step2.fill.fore_color.rgb = colors['light_gray']
    step2.line.color.rgb = colors['text']
    step2.line.width = Pt(1)
    
    step2_tf = step2.text_frame
    step2_tf.text = "② Sync to B360\n    triggers SLC logic"
    step2_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step2_tf.paragraphs[0].font.size = Pt(9)
    step2_tf.paragraphs[0].font.color.rgb = colors['text']
    step2_tf.paragraphs[0].font.name = "Arial"
    step2_tf.margin_left = Pt(8)
    step2_tf.margin_top = Pt(4)
    
    # Arrow
    flow_arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.9), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow2.fill.solid()
    flow_arrow2.fill.fore_color.rgb = colors['accent']
    flow_arrow2.line.fill.background()
    
    # Step 3 - highlighted
    step3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), step_y, step_width, step_height
    )
    step3.fill.solid()
    step3.fill.fore_color.rgb = RGBColor(255, 235, 235)
    step3.line.color.rgb = colors['accent']
    step3.line.width = Pt(2)
    
    step3_tf = step3.text_frame
    step3_tf.text = "③ B360 pushes back\n    (delete special code)"
    step3_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    step3_tf.paragraphs[0].font.size = Pt(9)
    step3_tf.paragraphs[0].font.bold = True
    step3_tf.paragraphs[0].font.color.rgb = colors['accent']
    step3_tf.paragraphs[0].font.name = "Arial"
    step3_tf.margin_left = Pt(8)
    step3_tf.margin_top = Pt(4)
    
    # Arrow
    flow_arrow3 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(7.25), Inches(4.38), arrow_width, Inches(0.22)
    )
    flow_arrow3.fill.solid()
    flow_arrow3.fill.fore_color.rgb = colors['accent']
    flow_arrow3.line.fill.background()
    
    # Step 4 - Problem result
    step4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.55), step_y, Inches(1.95), step_height
    )
    step4.fill.solid()
    step4.fill.fore_color.rgb = colors['accent']
    step4.line.fill.background()
    
    step4_tf = step4.text_frame
    step4_tf.text = "❌ Update fails\nor data overwritten"
    step4_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    step4_tf.paragraphs[0].font.size = Pt(9)
    step4_tf.paragraphs[0].font.bold = True
    step4_tf.paragraphs[0].font.color.rgb = colors['white']
    step4_tf.paragraphs[0].font.name = "Arial"
    
    # ========== Bottom: Solution ==========
    # Section title bar
    section5_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.95), Inches(0.06), Inches(0.28)
    )
    section5_bar.fill.solid()
    section5_bar.fill.fore_color.rgb = colors['success']
    section5_bar.line.fill.background()
    
    # Section title
    section5_title = slide.shapes.add_textbox(Inches(0.65), Inches(4.95), Inches(2.5), Inches(0.28))
    section5_frame = section5_title.text_frame
    section5_frame.text = "Proposed Solution"
    section5_frame.paragraphs[0].font.size = Pt(14)
    section5_frame.paragraphs[0].font.bold = True
    section5_frame.paragraphs[0].font.color.rgb = colors['success']
    section5_frame.paragraphs[0].font.name = "Arial"
    section5_frame.margin_left = 0
    section5_frame.margin_top = 0
    
    # Solution box
    solution_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.35), Inches(9.0), Inches(0.55)
    )
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    solution_box.line.color.rgb = colors['success']
    solution_box.line.width = Pt(2)
    
    # Solution content
    solution_items = [
        ("✓", "Migrate SLC logic to PMS side"),
        ("→", "PMS monitors reservation changes locally, decides special code updates"),
        ("→", "All changes initiated by PMS only"),
        ("→", "B360 acts as data receiver only"),
        ("→", "Avoid bi-directional override conflicts")
    ]
    
    solution_text = slide.shapes.add_textbox(Inches(0.6), Inches(5.35), Inches(8.8), Inches(0.55))
    solution_tf = solution_text.text_frame
    solution_tf.word_wrap = True
    
    for i, (symbol, text) in enumerate(solution_items):
        if i == 0:
            p = solution_tf.paragraphs[0]
        else:
            p = solution_tf.add_paragraph()
        p.text = f"{symbol} {text}"
        p.font.size = Pt(10)
        p.font.color.rgb = colors['text']
        p.font.name = "Arial"
        p.space_before = Pt(2)
        p.space_after = Pt(0)
        if i == 0:
            p.font.bold = True
    
    # Save file
    output_path = "/Users/emma/WorkBuddy/20260312092945/System_Sync_Issue.pptx"
    prs.save(output_path)
    print(f"\n✓ English PPT created successfully: {output_path}")
    return output_path

def main():
    """Main function"""
    print("=" * 60)
    print("System Data Synchronization Issue PPT - English Version")
    print("=" * 60)
    print()
    
    # Check and install dependencies
    if not check_and_install_dependencies():
        sys.exit(1)
    
    # Create presentation
    output_file = create_presentation()
    
    print("\n" + "=" * 60)
    print("Complete! English version PPT is ready")
    print("=" * 60)

if __name__ == "__main__":
    main()
